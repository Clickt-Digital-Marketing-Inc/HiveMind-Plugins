#!/usr/bin/env python3
"""
cm3_by_product.py — Clickt-branded per-product CM3 contribution report from a
Google Ads "Shopping products" CSV export.

Inputs:
  --csv <path>           Path to Google Ads Shopping products CSV
  --cogs-csv <path>      Optional Shopify "Gross profit by product" CSV. When
                         supplied, COGS% is taken per-product from this file
                         (with vendor + store-wide fallbacks). When omitted,
                         the blanket cogs_pct from --inputs is used.
  --inputs <path>        JSON with cogs_pct, ship_pct, proc_pct (and optional
                         fixed_costs, currency override). Defaults to the
                         previously used Clickt assumptions (65 / 20 / 2.9 / 0).
  --output-xlsx <path>   Output xlsx path (optional)
  --output-pptx <path>   Output pptx path (optional)
  --output-md   <path>   Output Obsidian-ready markdown path (optional)
  --period <str>         Optional period label (e.g. "Apr 1 – 30, 2026"). If
                         omitted, the script reads the date range from row 2
                         of the CSV.

If none of --output-xlsx, --output-pptx, --output-md is given, the script
defaults to writing all three with a timestamped filename in the current
working directory: clickt-cm3-by-product-YYYYMMDD-HHMMSS.{xlsx,pptx,md}

The script:
  1. Parses the CSV (skips title + date-range header rows).
  2. If a COGS CSV is provided, builds a title→COGS%, vendor→COGS%, and
     store-wide→COGS% lookup. Per-product COGS% is then resolved by:
       title match → vendor match (from ' : Vendor' suffix) → store-wide.
  3. Computes per-product CM1, CM2, CM3 using the resolved COGS% (or the
     blanket --inputs cogs_pct) plus the global ship_pct, proc_pct. Fixed
     costs (if supplied) are allocated to products by revenue share.
  4. Bands each ACTIVE product into 5 tiers based on CM3%:
     Excellent / High / Average / Low / Poor. Products with no spend and no
     revenue go to an Inactive bucket.
  5. Rolls up totals by Campaign, every Category level (L1–L5), every
     Product type level (L1–L5), and (if COGS CSV supplied) by Vendor.
  6. Writes any combination of Clickt-branded .xlsx, executive .pptx, and
     Obsidian-ready .md outputs.

Last stdout line is always a single JSON object listing the output paths and
headline KPIs — downstream callers (e.g. the Claude skill) read that.
"""

from __future__ import annotations

import argparse
import csv
import datetime as _dt
import json
import os
import re
import sys
from collections import defaultdict
from dataclasses import dataclass, field
from typing import Any

# Make Clickt brand module importable
PLUGIN_DIR = os.path.dirname(os.path.abspath(__file__))
if PLUGIN_DIR not in sys.path:
    sys.path.insert(0, PLUGIN_DIR)
import brand  # noqa: E402  (stdlib-safe: openpyxl is lazy-loaded inside brand)
import cm3_html  # noqa: E402  (stdlib-only renderer)

if False:  # typing only — never executed; keeps annotations resolvable to readers
    from openpyxl.worksheet.worksheet import Worksheet  # noqa: F401


# ─── LibreOffice xlsx normalization (open-reliably invariant) ─────────────────
def _find_soffice() -> str | None:
    """Locate the LibreOffice headless binary, or None."""
    import shutil
    if os.environ.get("CM3_FORCE_NO_SOFFICE") == "1":
        return None  # test hook: simulate a host without LibreOffice
    for name in ("soffice", "libreoffice"):
        p = shutil.which(name)
        if p:
            return p
    for p in (
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/bin/soffice",
        "/usr/local/bin/soffice",
    ):
        if os.path.exists(p):
            return p
    return None


def normalize_xlsx(path: str) -> None:
    """Round-trip an xlsx through LibreOffice so it opens reliably in Excel.

    Raises SystemExit(2) if soffice is unavailable — never silently ship an
    un-normalized file when normalization was requested.
    """
    import subprocess
    import tempfile

    soffice = _find_soffice()
    if not soffice:
        sys.stderr.write(
            "ERROR: LibreOffice (soffice) not found; cannot normalize the xlsx. "
            "Install LibreOffice or pass --no-normalize.\n"
        )
        raise SystemExit(2)
    with tempfile.TemporaryDirectory() as td:
        subprocess.run(
            [soffice, "--headless", "--calc", "--convert-to", "xlsx", "--outdir", td, path],
            check=True, capture_output=True,
        )
        produced = os.path.join(td, os.path.splitext(os.path.basename(path))[0] + ".xlsx")
        if not os.path.exists(produced):
            sys.stderr.write("ERROR: LibreOffice did not produce a normalized xlsx.\n")
            raise SystemExit(2)
        import shutil
        shutil.move(produced, path)


def check_workbook(path: str) -> int:
    """Integrity check for a built CM3 workbook. Exit 0 on pass, 1 on fail."""
    from openpyxl import load_workbook
    try:
        wb = load_workbook(path, data_only=True)
    except Exception as e:  # noqa: BLE001
        print(f"FAIL — cannot open {path}: {e}")
        return 1
    required = {"Summary", "By Product", "By Campaign", "Inputs & Methodology"}
    missing = required - set(wb.sheetnames)
    if missing:
        print(f"FAIL — missing sheets: {sorted(missing)}")
        return 1
    prod = wb["By Product"]
    if prod.max_row <= 4:
        print("FAIL — By Product tab has no product rows")
        return 1
    print(f"OK — {path}: {len(wb.sheetnames)} sheets, {prod.max_row - 4} product rows")
    return 0


# ─── Band thresholds (CM3% — share of revenue retained after var + ad + fixed) ─
# The four lower cutoffs (Excellent/High/Average/Low) are tunable: they default
# to the values below but can be overridden via --inputs or the --band-* CLI
# flags (and, live, by the in-Claude tuner). Poor is the implicit floor; band
# NAMES and styles never change. cutoffs_from()/band_thresholds() resolve the
# per-run values; classify() takes the resolved list so the tuned cutoffs flow
# through every output.
BAND_DEFS = [
    ("Excellent", "band_exc",  0.10, "clickt-band-strong"),
    ("High",      "band_high", 0.05, "clickt-band-healthy"),
    ("Average",   "band_avg",  0.00, "clickt-band-amber"),
    ("Low",       "band_low", -0.25, "clickt-band-amber"),
    ("Poor",      None,  float("-inf"), "clickt-band-red"),
]
BAND_NAMES = [b[0] for b in BAND_DEFS]
BAND_STYLE_NAME = {b[0]: b[3] for b in BAND_DEFS}
DEFAULT_CUTOFFS = {"band_exc": 0.10, "band_high": 0.05, "band_avg": 0.00, "band_low": -0.25}


def cutoffs_from(inputs: dict) -> dict:
    """Resolve the four tunable band cutoffs (fractions) from inputs, with defaults."""
    return {k: float(inputs.get(k, DEFAULT_CUTOFFS[k])) for k in DEFAULT_CUTOFFS}


def band_thresholds(inputs: dict) -> list[tuple[str, float, str]]:
    """Return [(name, lo, style), ...] with the tunable cutoffs applied."""
    c = cutoffs_from(inputs)
    out: list[tuple[str, float, str]] = []
    for name, key, default, style in BAND_DEFS:
        lo = c[key] if key else float("-inf")
        out.append((name, lo, style))
    return out


def _cutoff_pct(v: float) -> str:
    """0.10 -> '10%', 0.025 -> '2.5%', -0.25 -> '-25%', 0.0 -> '0%'."""
    s = f"{v * 100:.1f}".rstrip("0").rstrip(".")
    return (s or "0") + "%"


def band_range_texts(inputs: dict) -> dict:
    """Human-readable CM3% range per band, reflecting the tuned cutoffs."""
    c = cutoffs_from(inputs)
    e, h, a, l = (_cutoff_pct(c["band_exc"]), _cutoff_pct(c["band_high"]),
                  _cutoff_pct(c["band_avg"]), _cutoff_pct(c["band_low"]))
    return {
        "Excellent": f"CM3% ≥ {e}",
        "High":      f"{h} ≤ CM3% < {e}",
        "Average":   f"{a} ≤ CM3% < {h}",
        "Low":       f"{l} ≤ CM3% < {a}",
        "Poor":      f"CM3% < {l}  OR  spend with $0 revenue",
        "Inactive":  "$0 spend AND $0 revenue",
    }


def num(s: Any) -> float:
    """Parse a Google-Ads numeric cell. Handles '$', ',', '%', '--'."""
    if s is None:
        return 0.0
    if isinstance(s, (int, float)):
        return float(s)
    s = str(s).strip().replace(",", "").replace("$", "").replace("%", "")
    if s in ("", "--"):
        return 0.0
    try:
        return float(s)
    except ValueError:
        return 0.0


def clean_dim(s: str | None) -> str:
    """Normalize a dimension value. '--' or empty → ''."""
    if s is None:
        return ""
    s = s.strip()
    if s in ("", "--"):
        return ""
    return s


def classify(cm3_pct: float | None, revenue: float, ad_cost: float,
             thresholds: list[tuple[str, float, str]] | None = None) -> str:
    """Assign a band based on CM3%. Inactive if no rev + no spend.

    `thresholds` is the resolved band_thresholds(inputs) list; when omitted the
    default cutoffs are used (keeps ad-hoc callers working).
    """
    if revenue <= 0 and ad_cost <= 0:
        return "Inactive"
    if revenue <= 0:
        return "Poor"
    if cm3_pct is None:
        return "Poor"
    for name, lo, _ in (thresholds if thresholds is not None else band_thresholds({})):
        if cm3_pct >= lo:
            return name
    return "Poor"


@dataclass
class Product:
    title: str
    campaign: str
    cat: list[str]            # 5 levels, '' if absent
    ptype: list[str]          # 5 levels, '' if absent
    clicks: float
    impr: float
    cost: float               # ad spend
    conv: float
    conv_value: float         # revenue
    currency: str = "CAD"
    # COGS attribution
    vendor: str = ""          # extracted from " : Vendor" suffix
    cogs_pct: float = 0.0     # per-product, resolved
    cogs_source: str = ""     # "Title" / "Vendor" / "Store avg" / "Input"
    # computed
    cm1: float = 0.0
    cm2: float = 0.0
    fixed_alloc: float = 0.0
    cm3: float = 0.0
    cm3_pct: float | None = None
    band: str = ""


@dataclass
class Bucket:
    n_products: int = 0
    clicks: float = 0
    impr: float = 0
    cost: float = 0
    conv: float = 0
    conv_value: float = 0
    cm1: float = 0
    cm2: float = 0
    fixed_alloc: float = 0
    cm3: float = 0

    def add(self, p: Product) -> None:
        self.n_products += 1
        self.clicks += p.clicks
        self.impr += p.impr
        self.cost += p.cost
        self.conv += p.conv
        self.conv_value += p.conv_value
        self.cm1 += p.cm1
        self.cm2 += p.cm2
        self.fixed_alloc += p.fixed_alloc
        self.cm3 += p.cm3

    @property
    def cm3_pct(self) -> float | None:
        if self.conv_value <= 0:
            return None
        return self.cm3 / self.conv_value

    @property
    def roas(self) -> float | None:
        if self.cost <= 0:
            return None
        return self.conv_value / self.cost


# ─── Title normalization & vendor extraction (shared by GA ↔ GP matching) ────
_VENDOR_SUFFIX_RE = re.compile(r"\s*:\s*([^:]+?)(\s*\$[\d.,]+\s*[Pp]er\s+\w+\s*)?$")
_PRICE_SUFFIX_RE = re.compile(r"\s*\$[\d.,]+\s*[Pp]er\s+\w+\s*$")


def extract_vendor(title: str) -> str:
    """Pull the vendor token from a Google-Ads-style ' : Vendor' suffix."""
    if not title:
        return ""
    m = _VENDOR_SUFFIX_RE.search(title)
    if not m:
        return ""
    return m.group(1).strip()


def normalize_title(s: str) -> str:
    """Normalize a product title for cross-system matching."""
    s = s or ""
    s = _VENDOR_SUFFIX_RE.sub("", s)         # drop " : Vendor [$X.XX Per X]"
    s = _PRICE_SUFFIX_RE.sub("", s)          # drop a bare price suffix
    s = s.lower().strip()
    s = re.sub(r"\s+", " ", s)
    return s


# ─── COGS lookup from a Shopify "Gross profit by product" CSV ────────────────
@dataclass
class CogsLookup:
    by_title: dict[str, float] = field(default_factory=dict)
    by_vendor: dict[str, float] = field(default_factory=dict)
    store_avg: float | None = None
    # diagnostics
    n_products: int = 0
    n_vendors: int = 0
    total_sales: float = 0.0
    total_cogs: float = 0.0


def load_cogs_csv(path: str) -> CogsLookup:
    look = CogsLookup()
    by_vendor_acc: dict[str, list[tuple[float, float]]] = defaultdict(list)
    with open(path, "r", encoding="utf-8-sig") as f:
        for r in csv.DictReader(f):
            title = (r.get("Product title") or "").strip()
            vendor = (r.get("Product vendor") or "").strip()
            ns = num(r.get("Net sales"))
            cogs = num(r.get("Cost of goods sold"))
            if title and ns > 0:
                look.by_title[normalize_title(title)] = cogs / ns
                look.n_products += 1
                look.total_sales += ns
                look.total_cogs += cogs
                if vendor:
                    by_vendor_acc[vendor.lower().strip()].append((ns, cogs))
    for v, items in by_vendor_acc.items():
        ns = sum(x[0] for x in items)
        c = sum(x[1] for x in items)
        if ns > 0:
            look.by_vendor[v] = c / ns
    look.n_vendors = len(look.by_vendor)
    if look.total_sales > 0:
        look.store_avg = look.total_cogs / look.total_sales
    return look


def resolve_cogs_pct(p: Product, look: CogsLookup | None, fallback_pct: float) -> tuple[float, str]:
    """Return (cogs_pct, source_label)."""
    if look is None:
        return fallback_pct, "Input"
    n = normalize_title(p.title)
    if n in look.by_title:
        return look.by_title[n], "Title"
    v = p.vendor.lower().strip()
    if v and v in look.by_vendor:
        return look.by_vendor[v], "Vendor"
    if look.store_avg is not None:
        return look.store_avg, "Store avg"
    return fallback_pct, "Input"


# ─── CSV parsing ──────────────────────────────────────────────────────────────
def parse_csv(path: str) -> tuple[list[Product], str | None, str]:
    """Return (products, period_label, currency_code)."""
    with open(path, "r", encoding="utf-8-sig") as f:
        lines = f.readlines()

    # Header structure (Google Ads):
    #   line 0: "Shopping products"
    #   line 1: "<date range>"
    #   line 2: real column header
    period = None
    if len(lines) >= 2:
        period = lines[1].strip().strip('"')
    data_lines = lines[2:]
    reader = csv.DictReader(data_lines)

    products: list[Product] = []
    currency = "CAD"
    for r in reader:
        title = (r.get("Product Title") or "").strip()
        if not title:
            continue
        cur = clean_dim(r.get("Currency code"))
        if cur:
            currency = cur
        p = Product(
            title=title,
            campaign=clean_dim(r.get("Campaign")),
            cat=[clean_dim(r.get(f"Category ({_level_name(n)} level)")) for n in range(1, 6)],
            ptype=[clean_dim(r.get(f"Product type ({_level_name(n)} level)")) for n in range(1, 6)],
            clicks=num(r.get("Clicks")),
            impr=num(r.get("Impr.")),
            cost=num(r.get("Cost")),
            conv=num(r.get("Conversions")),
            conv_value=num(r.get("Conv. value")),
            currency=cur or currency,
            vendor=extract_vendor(title),
        )
        products.append(p)
    return products, period, currency


def _level_name(n: int) -> str:
    return {1: "1st", 2: "2nd", 3: "3rd", 4: "4th", 5: "5th"}[n]


# ─── Compute ──────────────────────────────────────────────────────────────────
def compute(products: list[Product], inputs: dict, cogs_look: CogsLookup | None = None) -> dict:
    fallback_cogs = inputs["cogs_pct"] / 100.0
    ship = inputs["ship_pct"] / 100.0
    proc = inputs["proc_pct"] / 100.0
    fixed = float(inputs.get("fixed_costs") or 0)

    thresholds = band_thresholds(inputs)
    total_rev = sum(p.conv_value for p in products)
    source_counts: dict[str, int] = defaultdict(int)
    source_revenue: dict[str, float] = defaultdict(float)
    for p in products:
        cogs_pct, src = resolve_cogs_pct(p, cogs_look, fallback_cogs)
        p.cogs_pct = cogs_pct
        p.cogs_source = src
        source_counts[src] += 1
        source_revenue[src] += p.conv_value
        var_pct = cogs_pct + ship + proc
        p.cm1 = p.conv_value * (1.0 - var_pct)
        p.cm2 = p.cm1 - p.cost
        if fixed > 0 and total_rev > 0:
            p.fixed_alloc = fixed * (p.conv_value / total_rev)
        else:
            p.fixed_alloc = 0.0
        p.cm3 = p.cm2 - p.fixed_alloc
        p.cm3_pct = (p.cm3 / p.conv_value) if p.conv_value > 0 else None
        p.band = classify(p.cm3_pct, p.conv_value, p.cost, thresholds)

    # Aggregate by band
    by_band: dict[str, Bucket] = {b: Bucket() for b in BAND_NAMES + ["Inactive"]}
    for p in products:
        by_band[p.band].add(p)

    # Aggregate by campaign
    by_campaign: dict[str, Bucket] = defaultdict(Bucket)
    for p in products:
        by_campaign[p.campaign or "(no campaign)"].add(p)

    # Aggregate by category and product type levels (only levels that have data)
    cat_levels = _aggregate_levels(products, key=lambda p: p.cat)
    pt_levels = _aggregate_levels(products, key=lambda p: p.ptype)

    # Aggregate by vendor (only if extracted)
    by_vendor: dict[str, Bucket] = defaultdict(Bucket)
    for p in products:
        if p.vendor:
            by_vendor[p.vendor].add(p)

    totals = Bucket()
    for p in products:
        totals.add(p)

    return dict(
        products=products,
        totals=totals,
        by_band=by_band,
        by_campaign=by_campaign,
        by_vendor=by_vendor,
        cat_levels=cat_levels,
        pt_levels=pt_levels,
        inputs=inputs,
        cutoffs=cutoffs_from(inputs),
        fallback_cogs=fallback_cogs, ship=ship, proc=proc, fixed=fixed,
        cogs_look=cogs_look,
        cogs_source_counts=dict(source_counts),
        cogs_source_revenue=dict(source_revenue),
    )


def _aggregate_levels(products: list[Product], key) -> list[tuple[int, dict[str, Bucket]]]:
    """Return [(level_n, {value: Bucket}), ...] only for levels with non-empty values."""
    out: list[tuple[int, dict[str, Bucket]]] = []
    for n in range(5):
        buckets: dict[str, Bucket] = defaultdict(Bucket)
        any_nonempty = False
        for p in products:
            v = key(p)[n]
            if v:
                any_nonempty = True
            buckets[v or "(unset)"].add(p)
        if any_nonempty:
            out.append((n + 1, dict(buckets)))
    return out


# ─── Writers ──────────────────────────────────────────────────────────────────
def _set_widths(ws: Worksheet, widths: list[tuple[str, int]]) -> None:
    for letter, w in widths:
        ws.column_dimensions[letter].width = w


def _h1(ws: Worksheet, row: int, text: str, span_cols: int = 8) -> None:
    cell = ws.cell(row=row, column=1, value=text)
    cell.style = "clickt-h1"
    ws.row_dimensions[row].height = 32
    if span_cols > 1:
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=span_cols)
        # Re-apply style to all cells in the merged range so the band fills
        for c in range(2, span_cols + 1):
            ws.cell(row=row, column=c).style = "clickt-h1"


def _h2(ws: Worksheet, row: int, text: str, span_cols: int = 8) -> None:
    cell = ws.cell(row=row, column=1, value=text)
    cell.style = "clickt-h2"
    ws.row_dimensions[row].height = 24
    if span_cols > 1:
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=span_cols)
        for c in range(2, span_cols + 1):
            ws.cell(row=row, column=c).style = "clickt-h2"


def _eyebrow(ws: Worksheet, row: int, text: str, span_cols: int = 8) -> None:
    cell = ws.cell(row=row, column=1, value=text)
    cell.style = "clickt-eyebrow"
    if span_cols > 1:
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=span_cols)


def _byline(ws: Worksheet, row: int, text: str, span_cols: int = 8) -> None:
    cell = ws.cell(row=row, column=1, value=text)
    cell.style = "clickt-byline"
    if span_cols > 1:
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=span_cols)


def _note(ws: Worksheet, row: int, text: str, span_cols: int = 8) -> None:
    cell = ws.cell(row=row, column=1, value=text)
    cell.style = "clickt-note"
    ws.row_dimensions[row].height = max(20, 14 * (1 + text.count("\n")))
    if span_cols > 1:
        ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=span_cols)


def _table_header(ws: Worksheet, row: int, headers: list[tuple[str, str]]) -> None:
    for i, (label, align) in enumerate(headers, start=1):
        cell = ws.cell(row=row, column=i, value=label)
        cell.style = "clickt-th-right" if align == "right" else "clickt-th"
    ws.row_dimensions[row].height = 24


def _write_band(ws: Worksheet, row: int, col: int, band: str) -> None:
    cell = ws.cell(row=row, column=col, value=band)
    if band == "Inactive":
        cell.style = "clickt-band-none"
    else:
        cell.style = BAND_STYLE_NAME.get(band, "clickt-band-none")


# ─── Sheet builders ───────────────────────────────────────────────────────────
def build_summary(ws: Worksheet, ctx: dict, period: str | None, currency: str) -> None:
    t: Bucket = ctx["totals"]
    by_band = ctx["by_band"]
    inp = ctx["inputs"]

    ws.title = "Summary"
    _set_widths(ws, [
        ("A", 28), ("B", 18), ("C", 18), ("D", 18), ("E", 18),
        ("F", 16), ("G", 16), ("H", 14),
    ])

    _eyebrow(ws, 1, "CLICKT  ·  GOOGLE ADS SHOPPING — CM3 BY PRODUCT")
    _h1(ws, 2, "Per-product contribution & band breakdown")
    _byline(ws, 3, f"Period: {period or 'unknown'}   ·   Currency: {currency}   ·   {len(ctx['products']):,} products")

    # Hero block
    ws.cell(row=5, column=1, value="Total CM3").style = "clickt-hero-label"
    ws.cell(row=5, column=2, value=t.cm3).style = "clickt-hero-currency"
    ws.cell(row=5, column=3, value="Total revenue").style = "clickt-hero-label"
    ws.cell(row=5, column=4, value=t.conv_value).style = "clickt-hero-currency"
    ws.cell(row=5, column=5, value="Total ad spend").style = "clickt-hero-label"
    ws.cell(row=5, column=6, value=t.cost).style = "clickt-hero-currency"
    ws.row_dimensions[5].height = 38

    # KPI table
    _h2(ws, 7, "Top-line KPIs", span_cols=8)
    _table_header(ws, 8, [
        ("Metric", "left"), ("Value", "right"), ("Formula", "left"),
        ("", "left"), ("", "left"), ("", "left"), ("", "left"), ("", "left"),
    ])

    def kpi_row(row, name, value, fmt_style, formula):
        ws.cell(row=row, column=1, value=name).style = "clickt-body"
        v_cell = ws.cell(row=row, column=2, value=value); v_cell.style = fmt_style
        ws.cell(row=row, column=3, value=formula).style = "clickt-body-dim"

    cm3_pct = (t.cm3 / t.conv_value) if t.conv_value > 0 else 0
    cm2_pct = (t.cm2 / t.conv_value) if t.conv_value > 0 else 0
    cm1_pct = (t.cm1 / t.conv_value) if t.conv_value > 0 else 0
    roas = (t.conv_value / t.cost) if t.cost > 0 else 0
    cpa = (t.cost / t.conv) if t.conv > 0 else 0
    aov = (t.conv_value / t.conv) if t.conv > 0 else 0
    ctr = (t.clicks / t.impr) if t.impr > 0 else 0
    cvr = (t.conv / t.clicks) if t.clicks > 0 else 0
    cpc = (t.cost / t.clicks) if t.clicks > 0 else 0

    kpi_row(9,  "Impressions",        t.impr,         "clickt-num-int",      "Sum of product-level Impr.")
    kpi_row(10, "Clicks",             t.clicks,       "clickt-num-int",      "Sum of product-level Clicks")
    kpi_row(11, "CTR",                ctr,            "clickt-num-pct",      "Clicks ÷ Impressions")
    kpi_row(12, "Avg. CPC",           cpc,            "clickt-num-currency-2","Ad spend ÷ Clicks")
    kpi_row(13, "Conversions",        t.conv,         "clickt-num-int",      "Sum of product-level Conversions")
    kpi_row(14, "Conversion rate",    cvr,            "clickt-num-pct",      "Conversions ÷ Clicks")
    kpi_row(15, "Avg. order value",   aov,            "clickt-num-currency-2","Conv. value ÷ Conversions")
    kpi_row(16, "Cost / conversion",  cpa,            "clickt-num-currency-2","Ad spend ÷ Conversions")
    kpi_row(17, "ROAS",               roas,           "clickt-num-multiple", "Conv. value ÷ Ad spend")
    kpi_row(18, "CM1",                t.cm1,          "clickt-num-currency", "Revenue × (1 − COGS% − Ship% − Proc%)")
    kpi_row(19, "CM1 %",              cm1_pct,        "clickt-num-pct",      "CM1 ÷ Revenue")
    kpi_row(20, "CM2",                t.cm2,          "clickt-num-currency", "CM1 − Ad spend")
    kpi_row(21, "CM2 %",              cm2_pct,        "clickt-num-pct",      "CM2 ÷ Revenue")
    if ctx["fixed"] > 0:
        kpi_row(22, "Fixed allocated", t.fixed_alloc, "clickt-num-currency", "Fixed × (Revenue ÷ Total revenue)")
        kpi_row(23, "CM3",            t.cm3,          "clickt-num-currency", "CM2 − Allocated fixed")
        kpi_row(24, "CM3 %",          cm3_pct,        "clickt-num-pct",      "CM3 ÷ Revenue")
        next_row = 26
    else:
        kpi_row(22, "CM3 (= CM2, no fixed alloc)", t.cm3, "clickt-num-currency", "CM2 (fixed costs not allocated)")
        kpi_row(23, "CM3 %",          cm3_pct,        "clickt-num-pct",      "CM3 ÷ Revenue")
        next_row = 25

    # Band summary
    _h2(ws, next_row, "Band breakdown (CM3% per product)", span_cols=8)
    _table_header(ws, next_row + 1, [
        ("Band", "left"), ("Products", "right"), ("Revenue", "right"),
        ("Ad spend", "right"), ("CM3", "right"), ("CM3 %", "right"),
        ("Share of CM3", "right"), ("", "left"),
    ])
    total_cm3 = t.cm3
    for i, band in enumerate(BAND_NAMES + ["Inactive"], start=next_row + 2):
        b: Bucket = by_band[band]
        _write_band(ws, i, 1, band)
        ws.cell(row=i, column=2, value=b.n_products).style = "clickt-num-int"
        ws.cell(row=i, column=3, value=b.conv_value).style = "clickt-num-currency"
        ws.cell(row=i, column=4, value=b.cost).style = "clickt-num-currency"
        ws.cell(row=i, column=5, value=b.cm3).style = "clickt-num-currency"
        cm3p = (b.cm3 / b.conv_value) if b.conv_value > 0 else 0
        ws.cell(row=i, column=6, value=cm3p).style = "clickt-num-pct"
        share = (b.cm3 / total_cm3) if total_cm3 != 0 else 0
        ws.cell(row=i, column=7, value=share).style = "clickt-num-pct"
        ws.cell(row=i, column=8, value="").style = "clickt-body-dim"

    band_end = next_row + 1 + len(BAND_NAMES) + 1
    _ct = cutoffs_from(ctx["inputs"])
    _band_line = (f"  Excellent ≥ {_cutoff_pct(_ct['band_exc'])}   ·   "
                  f"High {_cutoff_pct(_ct['band_high'])}–{_cutoff_pct(_ct['band_exc'])}   ·   "
                  f"Average {_cutoff_pct(_ct['band_avg'])}–{_cutoff_pct(_ct['band_high'])}   ·   "
                  f"Low {_cutoff_pct(_ct['band_low'])}–{_cutoff_pct(_ct['band_avg'])}   ·   "
                  f"Poor < {_cutoff_pct(_ct['band_low'])}")
    _note(ws, band_end + 2,
          "Band thresholds (CM3% = CM3 ÷ revenue):\n"
          + _band_line + "\n"
          "Products with $0 ad spend AND $0 revenue are 'Inactive' (excluded from active bands).\n"
          "Products with $0 revenue but >$0 ad spend are forced to Poor.",
          span_cols=8)


def build_band_detail(ws: Worksheet, ctx: dict) -> None:
    """One row per product, sorted by CM3 desc, with the band as a styled pill."""
    from openpyxl.utils import get_column_letter  # noqa: PLC0415  (lazy import)
    ws.title = "By Product"
    _set_widths(ws, [
        ("A", 56),  # title
        ("B", 22),  # vendor
        ("C", 22),  # cat L1
        ("D", 22),  # ptype L1
        ("E", 10),  # COGS%
        ("F", 11),  # COGS src
        ("G", 11),  # impr
        ("H", 10),  # clicks
        ("I", 12),  # cost
        ("J", 9),   # conv
        ("K", 12),  # conv val
        ("L", 9),   # roas
        ("M", 12),  # cm1
        ("N", 12),  # cm2
        ("O", 12),  # cm3
        ("P", 9),   # cm3%
        ("Q", 12),  # band
    ])

    _eyebrow(ws, 1, "DETAIL  ·  PRODUCTS, SORTED BY CM3 (DESC)")
    _h1(ws, 2, "By product", span_cols=17)
    _table_header(ws, 4, [
        ("Product title", "left"),
        ("Vendor", "left"),
        ("Cat L1", "left"),
        ("Product type L1", "left"),
        ("COGS %", "right"),
        ("COGS src", "left"),
        ("Impr.", "right"),
        ("Clicks", "right"),
        ("Ad spend", "right"),
        ("Conv.", "right"),
        ("Revenue", "right"),
        ("ROAS", "right"),
        ("CM1", "right"),
        ("CM2", "right"),
        ("CM3", "right"),
        ("CM3 %", "right"),
        ("Band", "left"),
    ])

    products = sorted(ctx["products"], key=lambda p: -p.cm3)
    for i, p in enumerate(products, start=5):
        ws.cell(row=i, column=1, value=p.title).style = "clickt-body"
        ws.cell(row=i, column=2, value=p.vendor).style = "clickt-body-dim"
        ws.cell(row=i, column=3, value=p.cat[0]).style = "clickt-body-dim"
        ws.cell(row=i, column=4, value=p.ptype[0]).style = "clickt-body-dim"
        ws.cell(row=i, column=5, value=p.cogs_pct).style = "clickt-num-pct"
        ws.cell(row=i, column=6, value=p.cogs_source).style = "clickt-body-dim"
        ws.cell(row=i, column=7, value=p.impr).style = "clickt-num-int"
        ws.cell(row=i, column=8, value=p.clicks).style = "clickt-num-int"
        ws.cell(row=i, column=9, value=p.cost).style = "clickt-num-currency-2"
        ws.cell(row=i, column=10, value=p.conv).style = "clickt-num-int"
        ws.cell(row=i, column=11, value=p.conv_value).style = "clickt-num-currency-2"
        roas = (p.conv_value / p.cost) if p.cost > 0 else None
        rcell = ws.cell(row=i, column=12, value=roas if roas is not None else "—")
        rcell.style = "clickt-num-multiple" if roas is not None else "clickt-body-dim"
        ws.cell(row=i, column=13, value=p.cm1).style = "clickt-num-currency-2"
        ws.cell(row=i, column=14, value=p.cm2).style = "clickt-num-currency-2"
        ws.cell(row=i, column=15, value=p.cm3).style = "clickt-num-currency-2"
        pct_cell = ws.cell(row=i, column=16, value=p.cm3_pct if p.cm3_pct is not None else "—")
        pct_cell.style = "clickt-num-pct" if p.cm3_pct is not None else "clickt-body-dim"
        _write_band(ws, i, 17, p.band)

    ws.freeze_panes = "A5"
    ws.auto_filter.ref = f"A4:{get_column_letter(17)}{4 + len(products)}"


def build_band_rollup(ws: Worksheet, ctx: dict) -> None:
    """One sheet per band, showing the products inside it (top 200 each)."""
    ws.title = "By Band"
    _set_widths(ws, [
        ("A", 56), ("B", 24), ("C", 22), ("D", 18), ("E", 12),
        ("F", 12), ("G", 12), ("H", 10),
    ])
    _eyebrow(ws, 1, "DETAIL  ·  PRODUCTS GROUPED BY BAND")
    _h1(ws, 2, "By band", span_cols=8)
    row = 4
    by_band = ctx["by_band"]
    products = ctx["products"]
    products_by_band: dict[str, list[Product]] = defaultdict(list)
    for p in products:
        products_by_band[p.band].append(p)
    for band in BAND_NAMES + ["Inactive"]:
        plist = sorted(products_by_band[band], key=lambda p: -p.cm3)
        b: Bucket = by_band[band]
        _h2(ws, row, f"{band}   ·   {b.n_products:,} products   ·   CM3 ${b.cm3:,.0f}   ·   Revenue ${b.conv_value:,.0f}",
            span_cols=8)
        row += 1
        _table_header(ws, row, [
            ("Product", "left"), ("Campaign", "left"), ("Cat L1", "left"),
            ("Product type L1", "left"), ("Ad spend", "right"),
            ("Revenue", "right"), ("CM3", "right"), ("CM3 %", "right"),
        ])
        row += 1
        cap = 150 if band in ("Low", "Poor", "Inactive") else 250
        for p in plist[:cap]:
            ws.cell(row=row, column=1, value=p.title).style = "clickt-body"
            ws.cell(row=row, column=2, value=p.campaign).style = "clickt-body-dim"
            ws.cell(row=row, column=3, value=p.cat[0]).style = "clickt-body-dim"
            ws.cell(row=row, column=4, value=p.ptype[0]).style = "clickt-body-dim"
            ws.cell(row=row, column=5, value=p.cost).style = "clickt-num-currency-2"
            ws.cell(row=row, column=6, value=p.conv_value).style = "clickt-num-currency-2"
            ws.cell(row=row, column=7, value=p.cm3).style = "clickt-num-currency-2"
            pct = p.cm3_pct
            pc = ws.cell(row=row, column=8, value=pct if pct is not None else "—")
            pc.style = "clickt-num-pct" if pct is not None else "clickt-body-dim"
            row += 1
        if len(plist) > cap:
            _note(ws, row, f"… {len(plist) - cap:,} more products in this band (see the By Product tab for the full list).", span_cols=8)
            row += 1
        row += 1


def build_campaign_rollup(ws: Worksheet, ctx: dict) -> None:
    ws.title = "By Campaign"
    _set_widths(ws, [
        ("A", 36), ("B", 12), ("C", 12), ("D", 12),
        ("E", 10), ("F", 12), ("G", 12), ("H", 12),
        ("I", 12), ("J", 10), ("K", 12),
    ])
    _eyebrow(ws, 1, "ROLLUP  ·  CAMPAIGNS")
    _h1(ws, 2, "By campaign", span_cols=11)
    _byline(ws, 3, "How is the ad account organized — and which campaigns earn vs. lose CM3?", span_cols=11)
    _table_header(ws, 5, [
        ("Campaign", "left"), ("Products", "right"), ("Impr.", "right"),
        ("Clicks", "right"), ("Conv.", "right"), ("Ad spend", "right"),
        ("Revenue", "right"), ("ROAS", "right"),
        ("CM3", "right"), ("CM3 %", "right"), ("Share of CM3", "right"),
    ])
    rows = sorted(ctx["by_campaign"].items(), key=lambda kv: -kv[1].cm3)
    total_cm3 = ctx["totals"].cm3
    for i, (name, b) in enumerate(rows, start=6):
        ws.cell(row=i, column=1, value=name).style = "clickt-body"
        ws.cell(row=i, column=2, value=b.n_products).style = "clickt-num-int"
        ws.cell(row=i, column=3, value=b.impr).style = "clickt-num-int"
        ws.cell(row=i, column=4, value=b.clicks).style = "clickt-num-int"
        ws.cell(row=i, column=5, value=b.conv).style = "clickt-num-int"
        ws.cell(row=i, column=6, value=b.cost).style = "clickt-num-currency-2"
        ws.cell(row=i, column=7, value=b.conv_value).style = "clickt-num-currency-2"
        roas = b.roas
        rc = ws.cell(row=i, column=8, value=roas if roas is not None else "—")
        rc.style = "clickt-num-multiple" if roas is not None else "clickt-body-dim"
        ws.cell(row=i, column=9, value=b.cm3).style = "clickt-num-currency-2"
        pct = b.cm3_pct
        pcl = ws.cell(row=i, column=10, value=pct if pct is not None else "—")
        pcl.style = "clickt-num-pct" if pct is not None else "clickt-body-dim"
        share = (b.cm3 / total_cm3) if total_cm3 else 0
        ws.cell(row=i, column=11, value=share).style = "clickt-num-pct"


def build_level_rollup(ws: Worksheet, ctx: dict, levels: list[tuple[int, dict[str, Bucket]]],
                       dim_label: str, sheet_title: str, eyebrow: str) -> None:
    ws.title = sheet_title
    _set_widths(ws, [
        ("A", 38), ("B", 12), ("C", 12), ("D", 12),
        ("E", 10), ("F", 12), ("G", 12), ("H", 12),
        ("I", 12), ("J", 10), ("K", 12),
    ])
    _eyebrow(ws, 1, eyebrow)
    _h1(ws, 2, sheet_title, span_cols=11)
    _byline(ws, 3,
            f"Rollup at each level of {dim_label.lower()}. Only levels with non-empty values are shown.",
            span_cols=11)
    row = 5
    total_cm3 = ctx["totals"].cm3
    if not levels:
        _note(ws, row, f"No {dim_label} data in this CSV.", span_cols=11)
        return
    for level_n, buckets in levels:
        _h2(ws, row, f"{dim_label} — Level {level_n}", span_cols=11)
        row += 1
        _table_header(ws, row, [
            (f"L{level_n} value", "left"), ("Products", "right"), ("Impr.", "right"),
            ("Clicks", "right"), ("Conv.", "right"), ("Ad spend", "right"),
            ("Revenue", "right"), ("ROAS", "right"),
            ("CM3", "right"), ("CM3 %", "right"), ("Share of CM3", "right"),
        ])
        row += 1
        sorted_items = sorted(buckets.items(), key=lambda kv: -kv[1].cm3)
        for name, b in sorted_items:
            ws.cell(row=row, column=1, value=name).style = "clickt-body"
            ws.cell(row=row, column=2, value=b.n_products).style = "clickt-num-int"
            ws.cell(row=row, column=3, value=b.impr).style = "clickt-num-int"
            ws.cell(row=row, column=4, value=b.clicks).style = "clickt-num-int"
            ws.cell(row=row, column=5, value=b.conv).style = "clickt-num-int"
            ws.cell(row=row, column=6, value=b.cost).style = "clickt-num-currency-2"
            ws.cell(row=row, column=7, value=b.conv_value).style = "clickt-num-currency-2"
            roas = b.roas
            rc = ws.cell(row=row, column=8, value=roas if roas is not None else "—")
            rc.style = "clickt-num-multiple" if roas is not None else "clickt-body-dim"
            ws.cell(row=row, column=9, value=b.cm3).style = "clickt-num-currency-2"
            pct = b.cm3_pct
            pcl = ws.cell(row=row, column=10, value=pct if pct is not None else "—")
            pcl.style = "clickt-num-pct" if pct is not None else "clickt-body-dim"
            share = (b.cm3 / total_cm3) if total_cm3 else 0
            ws.cell(row=row, column=11, value=share).style = "clickt-num-pct"
            row += 1
        row += 2


def build_vendor_rollup(ws: Worksheet, ctx: dict) -> None:
    ws.title = "By Vendor"
    _set_widths(ws, [
        ("A", 38), ("B", 12), ("C", 10), ("D", 12),
        ("E", 12), ("F", 12), ("G", 12), ("H", 10),
        ("I", 12), ("J", 10), ("K", 12),
    ])
    _eyebrow(ws, 1, "ROLLUP  ·  VENDOR (EXTRACTED FROM PRODUCT TITLE SUFFIX)")
    _h1(ws, 2, "By vendor", span_cols=11)
    _byline(ws, 3,
            "Vendor parsed from the ' : Vendor' suffix on the Google Ads product title.",
            span_cols=11)
    _table_header(ws, 5, [
        ("Vendor", "left"), ("Products", "right"), ("COGS %", "right"),
        ("Impr.", "right"), ("Clicks", "right"),
        ("Ad spend", "right"), ("Revenue", "right"), ("ROAS", "right"),
        ("CM3", "right"), ("CM3 %", "right"), ("Share of CM3", "right"),
    ])
    rows = sorted(ctx["by_vendor"].items(), key=lambda kv: -kv[1].cm3)
    total_cm3 = ctx["totals"].cm3
    look: CogsLookup | None = ctx.get("cogs_look")
    for i, (name, b) in enumerate(rows, start=6):
        ws.cell(row=i, column=1, value=name).style = "clickt-body"
        ws.cell(row=i, column=2, value=b.n_products).style = "clickt-num-int"
        # vendor COGS%
        vendor_cogs = None
        if look is not None:
            vendor_cogs = look.by_vendor.get(name.lower().strip())
        if vendor_cogs is None:
            vendor_cogs = ctx["fallback_cogs"]
        ws.cell(row=i, column=3, value=vendor_cogs).style = "clickt-num-pct"
        ws.cell(row=i, column=4, value=b.impr).style = "clickt-num-int"
        ws.cell(row=i, column=5, value=b.clicks).style = "clickt-num-int"
        ws.cell(row=i, column=6, value=b.cost).style = "clickt-num-currency-2"
        ws.cell(row=i, column=7, value=b.conv_value).style = "clickt-num-currency-2"
        roas = b.roas
        rc = ws.cell(row=i, column=8, value=roas if roas is not None else "—")
        rc.style = "clickt-num-multiple" if roas is not None else "clickt-body-dim"
        ws.cell(row=i, column=9, value=b.cm3).style = "clickt-num-currency-2"
        pct = b.cm3_pct
        pcl = ws.cell(row=i, column=10, value=pct if pct is not None else "—")
        pcl.style = "clickt-num-pct" if pct is not None else "clickt-body-dim"
        share = (b.cm3 / total_cm3) if total_cm3 else 0
        ws.cell(row=i, column=11, value=share).style = "clickt-num-pct"


def build_methodology(ws: Worksheet, ctx: dict, csv_path: str, period: str | None,
                      currency: str) -> None:
    ws.title = "Inputs & Methodology"
    _set_widths(ws, [("A", 30), ("B", 28), ("C", 60)])
    _eyebrow(ws, 1, "INPUTS  ·  ASSUMPTIONS & FORMULAS")
    _h1(ws, 2, "How this report is built", span_cols=3)

    _h2(ws, 4, "Inputs", span_cols=3)
    _table_header(ws, 5, [("Field", "left"), ("Value", "left"), ("Hint", "left")])
    inp = ctx["inputs"]
    look: CogsLookup | None = ctx.get("cogs_look")
    cogs_label = "Per-product (Shopify GP CSV)" if look else f"{inp['cogs_pct']}% (blanket)"
    rows = [
        ("Source CSV", os.path.basename(csv_path), "Google Ads → Reports → Shopping products → CSV"),
        ("Period",     period or "(unknown)", "Read from the second row of the CSV"),
        ("Currency",   currency, "Read from the 'Currency code' column"),
        ("COGS %",     cogs_label, "Per-product when a Shopify Gross-profit CSV is supplied; else flat fallback"),
        ("Delivery & shipping %", f"{inp['ship_pct']}%", "Fulfillment, freight, packaging — % of revenue"),
        ("Payment processing %", f"{inp['proc_pct']}%", "Shopify Payments, Stripe, etc."),
        ("Fixed costs ($/period)", f"${inp.get('fixed_costs', 0):,.2f}",
         "Total fixed costs for the period. Allocated to products by revenue share."),
    ]
    for i, (lbl, val, hint) in enumerate(rows, start=6):
        ws.cell(row=i, column=1, value=lbl).style = "clickt-body"
        ws.cell(row=i, column=2, value=val).style = "clickt-body-mono"
        ws.cell(row=i, column=3, value=hint).style = "clickt-body-dim"

    next_row = 6 + len(rows) + 2

    if look is not None:
        _h2(ws, next_row, "COGS resolution coverage", span_cols=3)
        next_row += 1
        _table_header(ws, next_row, [
            ("Source", "left"), ("Products", "right"), ("Revenue share", "right"),
        ])
        next_row += 1
        sc = ctx["cogs_source_counts"]
        sr = ctx["cogs_source_revenue"]
        tot_rev = sum(sr.values()) or 1.0
        for src in ("Title", "Vendor", "Store avg", "Input"):
            n = sc.get(src, 0)
            if n == 0 and src not in ("Title", "Vendor"):
                continue
            ws.cell(row=next_row, column=1, value=src).style = "clickt-body"
            ws.cell(row=next_row, column=2, value=n).style = "clickt-num-int"
            ws.cell(row=next_row, column=3, value=(sr.get(src, 0) / tot_rev)).style = "clickt-num-pct"
            next_row += 1
        _note(ws, next_row,
              f"Shopify GP CSV: {look.n_products:,} products, {look.n_vendors:,} vendors, "
              f"store avg COGS% = {(look.store_avg or 0)*100:.1f}%.\n"
              "Resolution order: exact title match → vendor (from ' : Vendor' suffix) → store-wide average → input fallback.",
              span_cols=3)
        next_row += 4

    _h2(ws, next_row, "Formulas", span_cols=3)
    next_row += 1
    _table_header(ws, next_row, [("Metric", "left"), ("Formula", "left"), ("Notes", "left")])
    next_row += 1
    formulas = [
        ("CM1", "Revenue × (1 − COGS% − Ship% − Proc%)",
         "Per-product gross contribution before ad spend."),
        ("CM2", "CM1 − Ad spend (Google Ads Cost)",
         "Per-product contribution after acquisition cost."),
        ("Allocated fixed", "Fixed × (Product revenue ÷ Total revenue)",
         "Zero if Fixed costs input is 0."),
        ("CM3", "CM2 − Allocated fixed",
         "If Fixed = 0, CM3 = CM2."),
        ("CM3 %", "CM3 ÷ Revenue",
         "Undefined when revenue = 0."),
        ("ROAS", "Conv. value ÷ Ad spend",
         "Same as Google Ads 'Conv. value / cost'."),
    ]
    for lbl, frm, note in formulas:
        ws.cell(row=next_row, column=1, value=lbl).style = "clickt-body"
        ws.cell(row=next_row, column=2, value=frm).style = "clickt-body-mono"
        ws.cell(row=next_row, column=3, value=note).style = "clickt-body-dim"
        next_row += 1

    next_row += 1
    _h2(ws, next_row, "Band thresholds", span_cols=3)
    next_row += 1
    _table_header(ws, next_row, [("Band", "left"), ("CM3% range", "left"), ("Meaning", "left")])
    next_row += 1
    _rng = band_range_texts(ctx["inputs"])
    bands = [
        ("Excellent", _rng["Excellent"],
         "Near the no-ads ceiling — almost no ad-attributed cost. Scale up."),
        ("High",      _rng["High"],
         "Profitable on ads. Keep and optimize."),
        ("Average",   _rng["Average"],
         "Thin but positive. Watch — small drift can flip to negative."),
        ("Low",       _rng["Low"],
         "Losing money. Pause unless strategic (new launch / brand)."),
        ("Poor",      _rng["Poor"],
         "Money pit. Pause immediately and review listing / targeting."),
        ("Inactive",  _rng["Inactive"],
         "No data — not in the bands. Listed separately."),
    ]
    for lbl, rng, mean in bands:
        _write_band(ws, next_row, 1, lbl)
        ws.cell(row=next_row, column=2, value=rng).style = "clickt-body-mono"
        ws.cell(row=next_row, column=3, value=mean).style = "clickt-body-dim"
        next_row += 1

    next_row += 1
    _note(ws, next_row,
          "Notes on shopping data structure:\n"
          " · This export is Pmax (Performance Max), so there are no Standard Shopping product groups — \n"
          "   asset groups / listing groups aren't exposed at row level. Aggregation here uses Campaign +\n"
          "   Google product Category (L1–L5) + merchant Product type (L1–L5).\n"
          " · The script consumes all five Category levels and all five Product type levels. Levels that\n"
          "   are empty for every row (e.g. Product type L2–L5 in this dataset) are auto-skipped from the\n"
          "   rollup so the report works for any merchant.\n"
          " · The same script will work on Standard Shopping exports — those just have multiple Campaigns\n"
          "   and the By Campaign rollup becomes more informative.",
          span_cols=3)
    next_row += 6

    _byline(ws, next_row, "Clickt — clickt.ca/tools", span_cols=3)


# ─── Top-level workbook builder (extracted from old main()) ───────────────────
def build_xlsx(ctx: dict, period: str | None, currency: str, csv_path: str, output_path: str) -> None:
    """Build the full Clickt xlsx workbook. Layout is byte-identical to the
    original inline logic in the source skill."""
    from openpyxl import Workbook  # noqa: PLC0415  (lazy — keeps module import dep-free)
    wb = Workbook()
    brand.register_named_styles(wb)
    # remove default sheet
    default_ws = wb.active
    wb.remove(default_ws)

    summary_ws = wb.create_sheet("Summary")
    build_summary(summary_ws, ctx, period, currency)

    band_ws = wb.create_sheet("By Band")
    build_band_rollup(band_ws, ctx)

    prod_ws = wb.create_sheet("By Product")
    build_band_detail(prod_ws, ctx)

    camp_ws = wb.create_sheet("By Campaign")
    build_campaign_rollup(camp_ws, ctx)

    if ctx["by_vendor"]:
        vend_ws = wb.create_sheet("By Vendor")
        build_vendor_rollup(vend_ws, ctx)

    cat_ws = wb.create_sheet("By Category")
    build_level_rollup(cat_ws, ctx, ctx["cat_levels"],
                       dim_label="Category",
                       sheet_title="By Category",
                       eyebrow="ROLLUP  ·  GOOGLE PRODUCT CATEGORY (L1–L5)")

    pt_ws = wb.create_sheet("By Product Type")
    build_level_rollup(pt_ws, ctx, ctx["pt_levels"],
                       dim_label="Product type",
                       sheet_title="By Product Type",
                       eyebrow="ROLLUP  ·  MERCHANT PRODUCT TYPE (L1–L5)")

    meth_ws = wb.create_sheet("Inputs & Methodology")
    build_methodology(meth_ws, ctx, csv_path, period, currency)

    wb.save(output_path)


# ─── PPTX writer ──────────────────────────────────────────────────────────────
def _hex_to_rgb(argb: str):
    """Convert an 8-char ARGB hex from brand.BRAND (e.g. 'FFF3B61C') to an
    (R, G, B) tuple — alpha dropped."""
    s = argb.lstrip("#")
    if len(s) == 8:
        s = s[2:]
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _pptx_rgb(argb: str):
    from pptx.dml.color import RGBColor  # local import — keeps --help dep-free
    r, g, b = _hex_to_rgb(argb)
    return RGBColor(r, g, b)


# Per-band chart bar colours (mapped to brand tokens)
_BAND_BAR_COLOR = {
    "Excellent": "FF2D7A4A",   # green
    "High":      "FF5BA89A",   # sage
    "Average":   "FFB8861B",   # amber
    "Low":       "FFD99A0A",   # yellow-deep
    "Poor":      "FFB33A28",   # red
}

FONT_SERIF = "Fraunces"
FONT_SANS  = "Inter"
FONT_MONO  = "JetBrains Mono"


def _fmt_money(v: float, currency: str = "USD") -> str:
    sign = "-" if v < 0 else ""
    return f"{sign}${abs(v):,.0f}"


def _fmt_money2(v: float) -> str:
    sign = "-" if v < 0 else ""
    return f"{sign}${abs(v):,.2f}"


def _fmt_pct(v: float | None) -> str:
    if v is None:
        return "—"
    return f"{v * 100:.1f}%"


def _fmt_mult(v: float | None) -> str:
    if v is None:
        return "—"
    return f"{v:.2f}×"


def _trunc(s: str, n: int = 60) -> str:
    s = s or ""
    return s if len(s) <= n else s[: n - 1] + "…"


def _pptx_set_bg(slide, rgb_argb: str) -> None:
    """Set the slide background fill to a solid RGB color."""
    from pptx.dml.color import RGBColor  # noqa: F401
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _pptx_rgb(rgb_argb)


def _pptx_add_textbox(slide, left, top, width, height, text, *,
                      font_name=FONT_SANS, size=14, bold=False,
                      color_argb="FF0B0F0E", align=None):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    # tf.text uses the existing paragraph
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _pptx_rgb(color_argb)
    return tb


def _pptx_yellow_band(slide, *, height_in=0.7):
    """Add the Clickt yellow header band across the top of the slide."""
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(height_in))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _pptx_rgb(brand.BRAND["yellow"])
    shape.line.fill.background()  # no border
    return shape


def _pptx_blank_slide(prs):
    # layout index 6 = blank in default Office theme
    return prs.slides.add_slide(prs.slide_layouts[6])


def build_pptx(ctx: dict, period: str | None, currency: str, output_path: str) -> None:
    """Write the 7-slide Clickt executive deck.

    The python-pptx import lives inside this function so `--help` works
    without python-pptx installed.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu  # noqa: F401
    from pptx.enum.text import PP_ALIGN  # noqa: F401
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    t: Bucket = ctx["totals"]
    by_band = ctx["by_band"]
    by_campaign = ctx["by_campaign"]
    products = ctx["products"]

    revenue = t.conv_value
    ad_spend = t.cost
    cm3 = t.cm3
    cm3_pct = (cm3 / revenue) if revenue > 0 else None
    roas = (revenue / ad_spend) if ad_spend > 0 else None
    mer = roas  # Single-channel paid spend — MER = ROAS by definition here.

    # ─── Slide 1 — Title ─────────────────────────────────────────────────────
    s1 = _pptx_blank_slide(prs)
    _pptx_yellow_band(s1, height_in=2.5)
    _pptx_add_textbox(s1, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4),
                      "CLICKT  ·  GOOGLE ADS SHOPPING — CM3 BY PRODUCT",
                      font_name=FONT_MONO, size=11, color_argb=brand.BRAND["ink"])
    _pptx_add_textbox(s1, Inches(0.6), Inches(1.0), Inches(12), Inches(1.2),
                      "CM3 by Product — Executive Summary",
                      font_name=FONT_SERIF, size=40, color_argb=brand.BRAND["ink"])
    _pptx_add_textbox(s1, Inches(0.6), Inches(2.7), Inches(12), Inches(0.5),
                      f"Period: {period or 'unknown'}   ·   Currency: {currency}   ·   {len(products):,} products",
                      font_name=FONT_SANS, size=16, color_argb=brand.BRAND["ink"])
    _pptx_add_textbox(s1, Inches(0.6), Inches(6.9), Inches(12), Inches(0.4),
                      "Clickt — clickt.ca/tools",
                      font_name=FONT_MONO, size=10, color_argb=brand.BRAND["fg_dim"])

    # ─── Slide 2 — Headline KPIs ─────────────────────────────────────────────
    s2 = _pptx_blank_slide(prs)
    _pptx_yellow_band(s2)
    _pptx_add_textbox(s2, Inches(0.6), Inches(0.12), Inches(12), Inches(0.5),
                      "Headline KPIs",
                      font_name=FONT_SERIF, size=24, color_argb=brand.BRAND["ink"])

    kpis = [
        ("Total revenue",  _fmt_money(revenue)),
        ("Total ad spend", _fmt_money(ad_spend)),
        ("CM3 $",          _fmt_money(cm3)),
        ("CM3 %",          _fmt_pct(cm3_pct)),
        ("ROAS",           _fmt_mult(roas)),
        ("MER",            _fmt_mult(mer)),
    ]
    top = 1.1
    row_h = 1.0
    for i, (lbl, val) in enumerate(kpis):
        y = top + i * row_h
        _pptx_add_textbox(s2, Inches(0.6), Inches(y), Inches(4.0), Inches(0.55),
                          lbl, font_name=FONT_MONO, size=14, bold=True,
                          color_argb=brand.BRAND["fg_dim"])
        _pptx_add_textbox(s2, Inches(4.6), Inches(y - 0.15), Inches(8.5), Inches(0.9),
                          val, font_name=FONT_SERIF, size=36,
                          color_argb=brand.BRAND["ink"], align="right")

    # ─── Slide 3 — CM3 band distribution chart ───────────────────────────────
    s3 = _pptx_blank_slide(prs)
    _pptx_yellow_band(s3)
    _pptx_add_textbox(s3, Inches(0.6), Inches(0.12), Inches(12), Inches(0.5),
                      "CM3 Band Distribution",
                      font_name=FONT_SERIF, size=24, color_argb=brand.BRAND["ink"])
    _pptx_add_textbox(s3, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
                      "Number of products in each CM3% band.",
                      font_name=FONT_SANS, size=13, color_argb=brand.BRAND["fg_dim"])

    chart_data = CategoryChartData()
    categories = ["Excellent", "High", "Average", "Low", "Poor"]
    chart_data.categories = categories
    chart_data.add_series("Products", [by_band[c].n_products for c in categories])

    chart_shape = s3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.6),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = False
    # Per-bar fill colours
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.vary_by_categories = True
    series = plot.series[0]
    # Setting per-point colours requires touching each data point
    for idx, cat in enumerate(categories):
        point = series.points[idx]
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = _pptx_rgb(_BAND_BAR_COLOR[cat])
        point.format.line.fill.background()

    # ─── Slide 4 — Top 10 by CM3$ ────────────────────────────────────────────
    s4 = _pptx_blank_slide(prs)
    _pptx_yellow_band(s4)
    _pptx_add_textbox(s4, Inches(0.6), Inches(0.12), Inches(12), Inches(0.5),
                      "Top 10 products by CM3 $",
                      font_name=FONT_SERIF, size=24, color_argb=brand.BRAND["ink"])
    top10 = sorted(products, key=lambda p: -p.cm3)[:10]
    _pptx_table(s4, top10, currency)

    # ─── Slide 5 — Bottom 10 by CM3$ ─────────────────────────────────────────
    s5 = _pptx_blank_slide(prs)
    _pptx_yellow_band(s5)
    _pptx_add_textbox(s5, Inches(0.6), Inches(0.12), Inches(12), Inches(0.5),
                      "Bottom 10 products by CM3 $ — loss leaders",
                      font_name=FONT_SERIF, size=24, color_argb=brand.BRAND["ink"])
    loss = sorted([p for p in products if p.band in ("Low", "Poor")], key=lambda p: p.cm3)[:10]
    if len(loss) < 10:
        # backfill with lowest-CM3 active products to keep table layout consistent
        active = sorted([p for p in products if p.band not in ("Inactive",) and p not in loss],
                        key=lambda p: p.cm3)
        loss = (loss + active)[:10]
    _pptx_table(s5, loss, currency)

    # ─── Slide 6 — Top 5 campaigns ───────────────────────────────────────────
    s6 = _pptx_blank_slide(prs)
    _pptx_yellow_band(s6)
    _pptx_add_textbox(s6, Inches(0.6), Inches(0.12), Inches(12), Inches(0.5),
                      "Top 5 campaigns by CM3 $",
                      font_name=FONT_SERIF, size=24, color_argb=brand.BRAND["ink"])
    camps = sorted(by_campaign.items(), key=lambda kv: -kv[1].cm3)[:5]
    _pptx_campaign_table(s6, camps)

    # ─── Slide 7 — What to do next ───────────────────────────────────────────
    s7 = _pptx_blank_slide(prs)
    _pptx_yellow_band(s7)
    _pptx_add_textbox(s7, Inches(0.6), Inches(0.12), Inches(12), Inches(0.5),
                      "What to do next",
                      font_name=FONT_SERIF, size=24, color_argb=brand.BRAND["ink"])

    bullets = _recommendation_bullets(ctx)
    by = 1.4
    for i, line in enumerate(bullets):
        _pptx_add_textbox(s7, Inches(0.7), Inches(by + i * 1.5), Inches(0.4), Inches(0.6),
                          "•", font_name=FONT_SERIF, size=36,
                          color_argb=brand.BRAND["yellow_deep"])
        _pptx_add_textbox(s7, Inches(1.2), Inches(by + i * 1.5 + 0.05), Inches(11.6), Inches(1.4),
                          line, font_name=FONT_SANS, size=18,
                          color_argb=brand.BRAND["ink"])

    prs.save(output_path)


def _pptx_table(slide, products: list[Product], currency: str) -> None:
    """5-column product table: Title · Revenue · CM3$ · CM3% · Band."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    n = len(products) or 1
    rows = n + 1  # +1 header
    cols = 5
    left = Inches(0.6)
    top = Inches(1.0)
    width = Inches(12.1)
    height = Inches(5.5)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    # Column widths
    tbl.columns[0].width = Inches(6.4)
    tbl.columns[1].width = Inches(1.6)
    tbl.columns[2].width = Inches(1.6)
    tbl.columns[3].width = Inches(1.1)
    tbl.columns[4].width = Inches(1.4)
    headers = ["Product", "Revenue", "CM3 $", "CM3 %", "Band"]
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _pptx_rgb(brand.BRAND["paper_warm"])
        tf = cell.text_frame
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if ci in (1, 2, 3) else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = h
        run.font.name = FONT_MONO
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = _pptx_rgb(brand.BRAND["fg_dim"])

    for ri, prod in enumerate(products, start=1):
        cells = [
            (_trunc(prod.title, 64),                FONT_SANS, "left",  brand.BRAND["ink"]),
            (_fmt_money(prod.conv_value),           FONT_MONO, "right", brand.BRAND["ink"]),
            (_fmt_money(prod.cm3),                  FONT_MONO, "right", brand.BRAND["ink"]),
            (_fmt_pct(prod.cm3_pct),                FONT_MONO, "right", brand.BRAND["ink"]),
            (prod.band,                             FONT_MONO, "left",  brand.BRAND["ink"]),
        ]
        for ci, (text, font, align, color) in enumerate(cells):
            cell = tbl.cell(ri, ci)
            tf = cell.text_frame
            tf.text = ""
            p = tf.paragraphs[0]
            p.alignment = (PP_ALIGN.RIGHT if align == "right"
                           else PP_ALIGN.CENTER if align == "center"
                           else PP_ALIGN.LEFT)
            run = p.add_run()
            run.text = text
            run.font.name = font
            run.font.size = Pt(11)
            run.font.color.rgb = _pptx_rgb(color)
        # Band pill fill on last cell
        band_color = _BAND_BAR_COLOR.get(prod.band)
        last = tbl.cell(ri, 4)
        if band_color:
            last.fill.solid()
            last.fill.fore_color.rgb = _pptx_rgb(band_color)
            # Switch band-cell text to paper for contrast
            for p in last.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = _pptx_rgb(brand.BRAND["paper"])


def _pptx_campaign_table(slide, camps: list[tuple[str, Bucket]]) -> None:
    """5-column campaign table: Campaign · Products · Revenue · CM3$ · CM3%."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    n = len(camps) or 1
    rows = n + 1
    cols = 5
    left = Inches(0.6)
    top = Inches(1.0)
    width = Inches(12.1)
    height = Inches(4.5)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(6.0)
    tbl.columns[1].width = Inches(1.3)
    tbl.columns[2].width = Inches(1.9)
    tbl.columns[3].width = Inches(1.7)
    tbl.columns[4].width = Inches(1.2)
    headers = ["Campaign", "Products", "Revenue", "CM3 $", "CM3 %"]
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _pptx_rgb(brand.BRAND["paper_warm"])
        tf = cell.text_frame
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if ci > 0 else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = h
        run.font.name = FONT_MONO
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = _pptx_rgb(brand.BRAND["fg_dim"])

    for ri, (name, b) in enumerate(camps, start=1):
        cells = [
            (_trunc(name, 70),               "left"),
            (f"{b.n_products:,}",            "right"),
            (_fmt_money(b.conv_value),       "right"),
            (_fmt_money(b.cm3),              "right"),
            (_fmt_pct(b.cm3_pct),            "right"),
        ]
        for ci, (text, align) in enumerate(cells):
            cell = tbl.cell(ri, ci)
            tf = cell.text_frame
            tf.text = ""
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if align == "right" else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = text
            run.font.name = FONT_MONO if ci > 0 else FONT_SANS
            run.font.size = Pt(12)
            run.font.color.rgb = _pptx_rgb(brand.BRAND["ink"])


def _recommendation_bullets(ctx: dict) -> list[str]:
    """3 deterministic action bullets driven by the data."""
    by_band = ctx["by_band"]
    by_campaign = ctx["by_campaign"]
    products = ctx["products"]
    counts = ctx.get("cogs_source_counts", {}) or {}

    poor: Bucket = by_band["Poor"]
    bullets: list[str] = []

    # 1) Pause N products in Poor band burning $X
    bullets.append(
        f"Pause {poor.n_products:,} products in the Poor band burning "
        f"{_fmt_money(poor.cost)} in ad spend per period."
    )

    # 2) Scale top Excellent products
    top_exc = sorted(
        [p for p in products if p.band == "Excellent"],
        key=lambda p: -p.cm3,
    )[:3]
    if top_exc:
        names = ", ".join(_trunc(p.title, 40) for p in top_exc)
        total_cm3 = sum(p.cm3 for p in top_exc)
        bullets.append(
            f"Scale the top {len(top_exc)} Excellent product"
            f"{'s' if len(top_exc) != 1 else ''} — "
            f"{names} — together generating {_fmt_money(total_cm3)} CM3."
        )
    else:
        bullets.append(
            "No products currently in the Excellent band — investigate why no SKU is hitting "
            f"CM3% ≥ {_cutoff_pct(cutoffs_from(ctx['inputs'])['band_exc'])}."
        )

    # 3) Either: COGS coverage warning, or: reinvest into top campaign
    if counts.get("Input", 0) > 0:
        n_unresolved = counts["Input"]
        bullets.append(
            f"Investigate COGS for {n_unresolved:,} unresolved products — the blanket fallback "
            "was applied because they couldn't be matched in the Shopify Gross-profit CSV."
        )
    else:
        camps = sorted(by_campaign.items(), key=lambda kv: -kv[1].cm3)
        if camps:
            name, b = camps[0]
            bullets.append(
                f"Reinvest top campaign CM3 — campaign \"{_trunc(name, 60)}\" drove "
                f"{_fmt_money(b.cm3)} CM3 at {_fmt_pct(b.cm3_pct)}."
            )
        else:
            bullets.append(
                "No campaign data found — confirm the Google Ads export contains a Campaign column."
            )

    return bullets


# ─── Markdown writer ──────────────────────────────────────────────────────────
def _md_escape(s: str) -> str:
    return (s or "").replace("|", r"\|").replace("\n", " ").strip()


def _md_row(cells: list[str]) -> str:
    return "| " + " | ".join(cells) + " |"


def _yaml_q(v) -> str:
    """Render v as a safe double-quoted YAML scalar (escapes \\ and ")."""
    s = str(v).replace("\\", "\\\\").replace('"', '\\"').replace("\n", " ")
    return f'"{s}"'


def build_markdown(ctx: dict, period: str | None, currency: str, output_path: str,
                   charts: bool = True) -> None:
    """Write the Obsidian-ready markdown deliverable.

    When charts is True (default), the declared Vega-Lite charts are rendered
    to static SVGs at the model's default params via the vendored chart module
    (cm3_html.render_static_charts — hard-fails with a clear message if
    vl-convert-python is missing; pass --no-charts to skip) into
    {stem}_charts/ next to the md, and referenced from a "## Charts" section.
    """
    t: Bucket = ctx["totals"]
    by_band = ctx["by_band"]
    by_campaign = ctx["by_campaign"]
    by_vendor = ctx["by_vendor"]
    products = ctx["products"]

    revenue = t.conv_value
    ad_spend = t.cost
    cm3 = t.cm3
    cm3_pct = (cm3 / revenue) if revenue > 0 else None
    roas = (revenue / ad_spend) if ad_spend > 0 else None
    mer = roas

    today = _dt.date.today().isoformat()
    inp = ctx["inputs"]
    cogs_mode = "per-product (Shopify GP CSV)" if ctx.get("cogs_look") else f"{inp['cogs_pct']}% blanket"

    # Static chart SVGs (generated, never authored) — written next to the md
    # as {stem}_charts/{id}.svg and referenced relatively so the md renders on
    # GitHub and in editor previews.
    chart_refs: list[tuple[str, str, str]] = []
    if charts:
        out_dir = os.path.dirname(os.path.abspath(output_path)) or "."
        stem = os.path.splitext(os.path.basename(output_path))[0]
        cdir = os.path.join(out_dir, f"{stem}_charts")
        os.makedirs(cdir, exist_ok=True)
        for cid, ctitle, svg in cm3_html.render_static_charts(ctx, only="md"):
            with open(os.path.join(cdir, f"{cid}.svg"), "w", encoding="utf-8") as f:
                f.write(svg)
            chart_refs.append((cid, ctitle, f"{stem}_charts/{cid}.svg"))

    out: list[str] = []

    # Frontmatter (provenance: period, currency, exact params, generated)
    out.append("---")
    out.append(f'title: {_yaml_q("CM3 by Product — " + (period or "report"))}')
    out.append(f"date: {today}")
    out.append(f'period: {_yaml_q(period or "unknown")}')
    out.append(f"currency: {currency}")
    out.append(f'cogs: {_yaml_q(cogs_mode)}')
    out.append(f"ship_pct: {inp['ship_pct']}")
    out.append(f"proc_pct: {inp['proc_pct']}")
    out.append(f"fixed_costs: {inp.get('fixed_costs', 0)}")
    _cut = ctx.get("cutoffs") or cutoffs_from(inp)
    out.append(f"band_exc: {_cut['band_exc']}")
    out.append(f"band_high: {_cut['band_high']}")
    out.append(f"band_avg: {_cut['band_avg']}")
    out.append(f"band_low: {_cut['band_low']}")
    out.append(f'generated: {_yaml_q(_dt.datetime.now().isoformat(timespec="seconds"))}')
    out.append("tags: [cm3, clickt, ecommerce, performance]")
    out.append("---")
    out.append("")

    # Heading + summary paragraph
    out.append(f"# CM3 by Product — {period or 'report'}")
    out.append("")
    out.append(
        f"This report covers **{len(products):,} products** for the period "
        f"**{period or 'unknown'}** (currency {currency}). "
        f"Total revenue was **{_fmt_money(revenue)}**, on **{_fmt_money(ad_spend)}** of ad spend, "
        f"producing **{_fmt_money(cm3)}** of CM3 "
        f"(**{_fmt_pct(cm3_pct)}** of revenue, **{_fmt_mult(roas)}** ROAS). "
        f"Band split — "
        f"Excellent: {by_band['Excellent'].n_products:,} · "
        f"High: {by_band['High'].n_products:,} · "
        f"Average: {by_band['Average'].n_products:,} · "
        f"Low: {by_band['Low'].n_products:,} · "
        f"Poor: {by_band['Poor'].n_products:,} · "
        f"Inactive: {by_band['Inactive'].n_products:,}."
    )
    out.append("")

    # KPI table
    out.append("## Headline KPIs")
    out.append("")
    out.append(_md_row(["Metric", "Value"]))
    out.append(_md_row(["---", "---:"]))
    out.append(_md_row(["Revenue",   _fmt_money(revenue)]))
    out.append(_md_row(["Ad spend",  _fmt_money(ad_spend)]))
    out.append(_md_row(["CM3 $",     _fmt_money(cm3)]))
    out.append(_md_row(["CM3 %",     _fmt_pct(cm3_pct)]))
    out.append(_md_row(["ROAS",      _fmt_mult(roas)]))
    out.append(_md_row(["MER",       _fmt_mult(mer)]))
    out.append("")

    # Band distribution
    out.append("## Band Distribution")
    out.append("")
    out.append(_md_row(["Band", "Products", "Revenue", "Ad spend", "CM3 $", "CM3 %", "Share of CM3"]))
    out.append(_md_row(["---", "---:", "---:", "---:", "---:", "---:", "---:"]))
    total_cm3 = cm3 if cm3 != 0 else 0.0
    for band in BAND_NAMES + ["Inactive"]:
        b: Bucket = by_band[band]
        share = (b.cm3 / total_cm3) if total_cm3 else None
        out.append(_md_row([
            band,
            f"{b.n_products:,}",
            _fmt_money(b.conv_value),
            _fmt_money(b.cost),
            _fmt_money(b.cm3),
            _fmt_pct(b.cm3_pct),
            _fmt_pct(share),
        ]))
    out.append("")

    # Charts (static SVGs at the default params; relative refs)
    if chart_refs:
        out.append("## Charts")
        out.append("")
        for _cid, ctitle, relpath in chart_refs:
            out.append(f"![{ctitle}]({relpath})")
            out.append("")

    # Top 10 by CM3$
    out.append("## Top 10 by CM3 $")
    out.append("")
    out.append(_md_row(["Product", "Revenue", "CM3 $", "CM3 %", "Band"]))
    out.append(_md_row(["---", "---:", "---:", "---:", "---"]))
    for p in sorted(products, key=lambda p: -p.cm3)[:10]:
        out.append(_md_row([
            _md_escape(_trunc(p.title, 80)),
            _fmt_money(p.conv_value),
            _fmt_money(p.cm3),
            _fmt_pct(p.cm3_pct),
            p.band,
        ]))
    out.append("")

    # Bottom 10 by CM3$
    out.append("## Bottom 10 by CM3 $")
    out.append("")
    out.append(_md_row(["Product", "Revenue", "CM3 $", "CM3 %", "Band"]))
    out.append(_md_row(["---", "---:", "---:", "---:", "---"]))
    loss = sorted([p for p in products if p.band in ("Low", "Poor")], key=lambda p: p.cm3)[:10]
    if len(loss) < 10:
        active = sorted([p for p in products if p.band != "Inactive" and p not in loss],
                        key=lambda p: p.cm3)
        loss = (loss + active)[:10]
    for p in loss:
        out.append(_md_row([
            _md_escape(_trunc(p.title, 80)),
            _fmt_money(p.conv_value),
            _fmt_money(p.cm3),
            _fmt_pct(p.cm3_pct),
            p.band,
        ]))
    out.append("")

    # Rollups
    out.append("## Rollups")
    out.append("")

    # By Campaign
    out.append("### By Campaign")
    out.append("")
    out.append(_md_row(["Campaign", "Products", "Revenue", "Ad spend", "CM3 $", "CM3 %", "Share of CM3"]))
    out.append(_md_row(["---", "---:", "---:", "---:", "---:", "---:", "---:"]))
    for name, b in sorted(by_campaign.items(), key=lambda kv: -kv[1].cm3)[:10]:
        share = (b.cm3 / total_cm3) if total_cm3 else None
        out.append(_md_row([
            _md_escape(name),
            f"{b.n_products:,}",
            _fmt_money(b.conv_value),
            _fmt_money(b.cost),
            _fmt_money(b.cm3),
            _fmt_pct(b.cm3_pct),
            _fmt_pct(share),
        ]))
    out.append("")

    # By Category — L1 (and lower levels if present)
    cat_levels = ctx.get("cat_levels") or []
    if cat_levels:
        for level_n, buckets in cat_levels[:1]:  # L1 only — keep markdown compact
            out.append(f"### By Category — L{level_n}")
            out.append("")
            out.append(_md_row(["Category", "Products", "Revenue", "Ad spend", "CM3 $", "CM3 %"]))
            out.append(_md_row(["---", "---:", "---:", "---:", "---:", "---:"]))
            for name, b in sorted(buckets.items(), key=lambda kv: -kv[1].cm3):
                out.append(_md_row([
                    _md_escape(name),
                    f"{b.n_products:,}",
                    _fmt_money(b.conv_value),
                    _fmt_money(b.cost),
                    _fmt_money(b.cm3),
                    _fmt_pct(b.cm3_pct),
                ]))
            out.append("")

    # By Product Type — L1
    pt_levels = ctx.get("pt_levels") or []
    if pt_levels:
        for level_n, buckets in pt_levels[:1]:
            out.append(f"### By Product Type — L{level_n}")
            out.append("")
            out.append(_md_row(["Product Type", "Products", "Revenue", "Ad spend", "CM3 $", "CM3 %"]))
            out.append(_md_row(["---", "---:", "---:", "---:", "---:", "---:"]))
            for name, b in sorted(buckets.items(), key=lambda kv: -kv[1].cm3):
                out.append(_md_row([
                    _md_escape(name),
                    f"{b.n_products:,}",
                    _fmt_money(b.conv_value),
                    _fmt_money(b.cost),
                    _fmt_money(b.cm3),
                    _fmt_pct(b.cm3_pct),
                ]))
            out.append("")

    # By Vendor
    if by_vendor:
        out.append("### By Vendor")
        out.append("")
        out.append(_md_row(["Vendor", "Products", "Revenue", "Ad spend", "CM3 $", "CM3 %"]))
        out.append(_md_row(["---", "---:", "---:", "---:", "---:", "---:"]))
        for name, b in sorted(by_vendor.items(), key=lambda kv: -kv[1].cm3):
            out.append(_md_row([
                _md_escape(name),
                f"{b.n_products:,}",
                _fmt_money(b.conv_value),
                _fmt_money(b.cost),
                _fmt_money(b.cm3),
                _fmt_pct(b.cm3_pct),
            ]))
        out.append("")

    # All products — full no-row-loss table (every input row carries a band status)
    out.append("## All products (full)")
    out.append("")
    out.append(f"Every one of the {len(products):,} products, sorted by CM3 $ — no rows dropped.")
    out.append("")
    out.append(_md_row(["Product", "Campaign", "Vendor", "Ad spend", "Revenue", "CM3 $", "CM3 %", "Band"]))
    out.append(_md_row(["---", "---", "---", "---:", "---:", "---:", "---:", "---"]))
    for p in sorted(products, key=lambda p: -p.cm3):
        out.append(_md_row([
            _md_escape(p.title),
            _md_escape(p.campaign or "(no campaign)"),
            _md_escape(p.vendor),
            _fmt_money(p.cost),
            _fmt_money(p.conv_value),
            _fmt_money(p.cm3),
            _fmt_pct(p.cm3_pct),
            p.band,
        ]))
    out.append("")

    # Actions w/ callouts
    out.append("## What to do next")
    out.append("")
    bullets = _recommendation_bullets(ctx)
    # Bullet 1 → warning (Poor band)
    out.append("> [!warning]")
    out.append(f"> {bullets[0]}")
    out.append("")
    # Bullet 2 → success (Excellent scale-up)
    out.append("> [!success]")
    out.append(f"> {bullets[1]}")
    out.append("")
    # Bullet 3 → plain bullet
    out.append(f"- {bullets[2]}")
    out.append("")

    # Related
    out.append("## Related")
    out.append("")
    out.append("- [[CM3 Calculator]]")
    out.append("- [[Max CAC]]")
    out.append("")

    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(out))


# ─── Main ─────────────────────────────────────────────────────────────────────
def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(
        description="Per-product CM3 contribution report — writes any combination of xlsx, pptx, and md."
    )
    ap.add_argument("--csv", default=None, help="Path to the Google Ads Shopping products CSV.")
    ap.add_argument("--cogs-csv", dest="cogs_csv", default=None,
                    help="Optional Shopify 'Gross profit by product' CSV for per-product COGS.")
    ap.add_argument("--inputs", default=None,
                    help="Path to cm3-by-product-inputs.json (defaults to bundled assumptions).")
    ap.add_argument("--output-xlsx", dest="output_xlsx", default=None,
                    help="Output path for the detailed Clickt-branded xlsx workbook.")
    ap.add_argument("--output-html", dest="output_html", default=None,
                    help="Output path for the self-contained interactive HTML explorer.")
    ap.add_argument("--output-md", dest="output_md", default=None,
                    help="Output path for the Obsidian-ready markdown doc.")
    ap.add_argument("--output-pptx", dest="output_pptx", default=None,
                    help="Output path for the executive PowerPoint deck (opt-in; not in the default bundle).")
    ap.add_argument("--pptx", action="store_true",
                    help="In default mode, also emit the 7-slide PowerPoint deck (off by default).")
    ap.add_argument("--no-normalize", dest="normalize", action="store_false",
                    help="Skip LibreOffice normalization of the xlsx (normalized by default).")
    ap.set_defaults(normalize=True)
    ap.add_argument("--no-charts", dest="charts", action="store_false",
                    help="Skip the generated Vega-Lite charts (static md SVGs, the live "
                         "explorer charts, and the tuner widget chart). Charts are on by "
                         "default and require vl-convert-python for the static renders.")
    ap.set_defaults(charts=True)
    ap.add_argument("--check", default=None, metavar="XLSX",
                    help="Integrity-check a built workbook and exit (no CSV needed).")
    ap.add_argument("--period", default=None,
                    help="Optional period override (e.g. 'Apr 1 – 30, 2026').")
    ap.add_argument("--brand", default=None,
                    help="Operator/brand label used in the tuner + Save/Export prompts.")
    ap.add_argument("--emit-widget", dest="emit_widget", default=None,
                    help="Write the in-Claude tuner HTML fragment to this path and exit "
                         "(no md/html/xlsx/pptx are built).")
    # Tunable parameter overrides (take precedence over --inputs / defaults). These
    # let a tuner Save/Export be a single one-shot command at the operator's settings.
    ap.add_argument("--cogs-pct", dest="cogs_pct", type=float, default=None,
                    help="Blanket COGS %% fallback (used when no per-product Shopify COGS).")
    ap.add_argument("--ship-pct", dest="ship_pct", type=float, default=None, help="Shipping %%.")
    ap.add_argument("--proc-pct", dest="proc_pct", type=float, default=None, help="Payment processing %%.")
    ap.add_argument("--fixed-costs", dest="fixed_costs", type=float, default=None,
                    help="Fixed costs (allocated pro-rata by revenue).")
    ap.add_argument("--band-exc", dest="band_exc", type=float, default=None,
                    help="Excellent band cutoff (CM3%% fraction, e.g. 0.10).")
    ap.add_argument("--band-high", dest="band_high", type=float, default=None,
                    help="High band cutoff (fraction, e.g. 0.05).")
    ap.add_argument("--band-avg", dest="band_avg", type=float, default=None,
                    help="Average band cutoff (fraction, e.g. 0.00).")
    ap.add_argument("--band-low", dest="band_low", type=float, default=None,
                    help="Low band cutoff (fraction, e.g. -0.25).")
    args = ap.parse_args(argv)

    if args.check:
        return check_workbook(args.check)

    if not args.csv:
        ap.error("--csv is required (unless using --check)")

    # Default bundle = md + html + xlsx (pptx only when --pptx). Explicit
    # --output-* paths override and are always honored. --emit-widget never
    # triggers the default bundle (it builds only the tuner fragment).
    explicit = any([args.output_xlsx, args.output_html, args.output_md, args.output_pptx])
    if not explicit and not args.emit_widget:
        ts = _dt.datetime.now().strftime("%Y%m%d-%H%M%S")
        base = os.path.join(os.getcwd(), f"clickt-cm3-by-product-{ts}")
        args.output_md = base + ".md"
        args.output_html = base + "_explorer.html"
        args.output_xlsx = base + ".xlsx"
        if args.pptx:
            args.output_pptx = base + ".pptx"

    # Load inputs
    if args.inputs and os.path.exists(args.inputs):
        with open(args.inputs, "r") as f:
            inputs = json.load(f)
    else:
        inputs = {}
    inputs.setdefault("cogs_pct", 65)
    inputs.setdefault("ship_pct", 20)
    inputs.setdefault("proc_pct", 2.9)
    inputs.setdefault("fixed_costs", 0)
    for _k in DEFAULT_CUTOFFS:
        inputs.setdefault(_k, DEFAULT_CUTOFFS[_k])
    # CLI overrides win over --inputs / defaults (one-shot tuner Save/Export).
    for _k in ("cogs_pct", "ship_pct", "proc_pct", "fixed_costs",
               "band_exc", "band_high", "band_avg", "band_low"):
        _v = getattr(args, _k, None)
        if _v is not None:
            inputs[_k] = _v

    products, period, currency = parse_csv(args.csv)
    if args.period:
        period = args.period

    cogs_look = None
    if args.cogs_csv and os.path.exists(args.cogs_csv):
        cogs_look = load_cogs_csv(args.cogs_csv)

    ctx = compute(products, inputs, cogs_look=cogs_look)

    # In-Claude tuner fragment (no report files built). The widget's Save/Export
    # buttons reproduce md/xlsx/html/pptx at the tuned params via this same CLI.
    if args.emit_widget:
        frag = cm3_html.build_widget_fragment(
            ctx, period, currency,
            brand=args.brand, csv_path=args.csv, cogs_csv=args.cogs_csv,
            charts=args.charts,
        )
        with open(args.emit_widget, "w", encoding="utf-8") as f:
            f.write(frag)
        print(f"Wrote widget {args.emit_widget}")
        return 0

    written: dict[str, str] = {}
    if args.output_md:
        build_markdown(ctx, period, currency, args.output_md, charts=args.charts)
        written["md"] = os.path.abspath(args.output_md)
        print(f"Wrote {args.output_md}")
    if args.output_html:
        cm3_html.build_html(ctx, period, currency, args.output_html, charts=args.charts)
        written["html"] = os.path.abspath(args.output_html)
        print(f"Wrote {args.output_html}")
    if args.output_xlsx:
        build_xlsx(ctx, period, currency, args.csv, args.output_xlsx)
        if args.normalize:
            normalize_xlsx(args.output_xlsx)
        written["xlsx"] = os.path.abspath(args.output_xlsx)
        print(f"Wrote {args.output_xlsx}")
    if args.output_pptx:
        build_pptx(ctx, period, currency, args.output_pptx)
        written["pptx"] = os.path.abspath(args.output_pptx)
        print(f"Wrote {args.output_pptx}")

    # human-readable summary line (kept for backward-compat with the original
    # source's stdout — but the JSON line below is the authoritative one)
    t: Bucket = ctx["totals"]
    bands = ctx["by_band"]
    parts = [f"{b}={bands[b].n_products}" for b in BAND_NAMES + ["Inactive"]]
    cm3_pct = (t.cm3 / t.conv_value) if t.conv_value > 0 else 0
    roas = (t.conv_value / t.cost) if t.cost > 0 else 0
    print(f"  {len(products):,} products · Revenue ${t.conv_value:,.0f} · Ad spend ${t.cost:,.0f} · "
          f"ROAS {roas:.2f}× · CM3 ${t.cm3:,.0f} ({cm3_pct*100:.1f}%)")
    print("  Bands: " + " · ".join(parts))

    # Final stdout line — single JSON object, machine-readable.
    summary = {
        **written,
        "revenue":         t.conv_value,
        "ad_spend":        t.cost,
        "cm3":             t.cm3,
        "cm3_pct":         cm3_pct,
        "roas":            roas,
        "excellent_count": bands["Excellent"].n_products,
        "high_count":      bands["High"].n_products,
        "average_count":   bands["Average"].n_products,
        "low_count":       bands["Low"].n_products,
        "poor_count":      bands["Poor"].n_products,
        "inactive_count":  bands["Inactive"].n_products,
        "period":          period or "",
        "currency":        currency,
    }
    print(json.dumps(summary))
    return 0


if __name__ == "__main__":
    sys.exit(main())
